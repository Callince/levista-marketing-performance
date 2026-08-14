"""Levista_Performance_Report.pptx — the 21-slide management deck.

Charts are native PowerPoint charts (not images) so the team can restyle or
re-point them without regenerating the deck.
"""
from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

import config
from analytics import insights as ai
from analytics import metrics
from db.models import alerts, get_engine, insights as insights_table, recommendations

BRAND = RGBColor(0x1F, 0x3B, 0x4D)
ACCENT = RGBColor(0x6F, 0x4E, 0x37)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREY = RGBColor(0x55, 0x60, 0x6A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF5, 0xF7)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)

# The brand wordmark is dark brown; slide furniture is dark navy, so the white
# variant is the one that reads. Both are generated from the source SVG by
# tools/build_logo.py.
ASSETS = Path(__file__).resolve().parent.parent / "assets"
LOGO_WHITE = ASSETS / "levista_logo_white.png"
LOGO_RATIO = 3713 / 1073          # from the trimmed raster


def _text(frame, lines, size=14, color=None, bold_first=False, space=6):
    frame.word_wrap = True
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = str(line)
        para.space_after = Pt(space)
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color or GREY
            run.font.bold = bold_first and i == 0
    return frame


def _slide(prs, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()
    frame = bar.text_frame
    frame.margin_left = MARGIN
    para = frame.paragraphs[0]
    para.text = title
    para.runs[0].font.size = Pt(26)
    para.runs[0].font.bold = True
    para.runs[0].font.color.rgb = WHITE
    if LOGO_WHITE.exists():
        height = Inches(0.34)
        width = int(height * LOGO_RATIO)
        slide.shapes.add_picture(str(LOGO_WHITE), W - MARGIN - width,
                                 int((Inches(0.95) - height) / 2), width, height)
    if subtitle:
        sub = slide.shapes.add_textbox(MARGIN, Inches(1.0), W - 2 * MARGIN, Inches(0.4))
        _text(sub.text_frame, [subtitle], size=12, color=GREY)
    return slide


def _kpi_cards(slide, cards, top=Inches(1.5), height=Inches(1.15)):
    count = len(cards)
    gap = Inches(0.15)
    width = int((W - 2 * MARGIN - gap * (count - 1)) / count)
    for i, (label, value, tone) in enumerate(cards):
        left = MARGIN + i * (width + gap)
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDC)
        frame = box.text_frame
        frame.word_wrap = True
        head = frame.paragraphs[0]
        head.text = label
        head.alignment = PP_ALIGN.CENTER
        head.runs[0].font.size = Pt(11)
        head.runs[0].font.color.rgb = GREY
        body = frame.add_paragraph()
        body.text = value
        body.alignment = PP_ALIGN.CENTER
        body.runs[0].font.size = Pt(22)
        body.runs[0].font.bold = True
        body.runs[0].font.color.rgb = tone or BRAND
    return top + height + Inches(0.25)


def _chart(slide, kind, categories, series: dict, left, top, width, height,
           title=None, legend=True):
    categories = list(categories)
    # python-pptx raises "chart data contains no categories" on an empty series, which
    # takes the whole deck — and therefore the whole rebuild job — down. An empty
    # slice is a normal state (a fresh install, one platform, a filtered view), so
    # every chart degrades to a note instead of failing.
    if not categories:
        box = slide.shapes.add_textbox(left, top, width, height)
        _text(box.text_frame, [title or "No data", "Nothing to chart for this selection."],
              size=12, bold_first=True, color=GREY)
        return None

    data = CategoryChartData()
    data.categories = categories
    for name, values in series.items():
        data.add_series(name, [None if v is None or pd.isna(v) else float(v) for v in values])
    graphic = slide.shapes.add_chart(kind, left, top, width, height, data)
    chart = graphic.chart
    chart.has_title = title is not None
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    chart.has_legend = legend and len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    try:
        chart.font.size = Pt(10)
    except Exception:
        pass
    return chart


def _table(slide, df, columns, left, top, width, height, formatters=None):
    formatters = formatters or {}
    rows = len(df) + 1
    shape = slide.shapes.add_table(rows, len(columns), left, top, width, height)
    table = shape.table
    for c, (header, _) in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(11)
        para.runs[0].font.bold = True
        para.runs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND
    for r, (_, record) in enumerate(df.iterrows(), start=1):
        for c, (_, field) in enumerate(columns):
            value = record.get(field)
            text = formatters.get(field, lambda v: "—" if v is None or pd.isna(v) else str(v))(value)
            cell = table.cell(r, c)
            cell.text = text
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(10)
            if field == "roas" and pd.notna(value):
                para.runs[0].font.color.rgb = (
                    GREEN if value >= config.HEALTHY_ROAS
                    else AMBER if value >= config.BREAKEVEN_ROAS else RED)
                para.runs[0].font.bold = True
    return shape


# ---------------------------------------------------------------- formatters
def f_money(v):
    return ai.inr(v)


def f_int(v):
    return "—" if v is None or pd.isna(v) else f"{int(v):,}"


def f_roas(v):
    return "—" if v is None or pd.isna(v) else f"{v:.2f}"


def f_pct(v):
    return ai.pct(v)


def f_text(v):
    return "—" if v is None or (isinstance(v, float) and pd.isna(v)) else str(v)[:46]


FMT = {"spend": f_money, "revenue": f_money, "impact_value": f_money, "value": f_money,
       "orders": f_int, "units": f_int, "impressions": f_int, "clicks": f_int,
       "roas": f_roas, "ctr": f_pct, "conv_rate": f_pct, "revenue_share": f_pct,
       "campaign_name": f_text, "product_name": f_text, "keyword": f_text,
       "city": f_text, "entity_name": f_text, "platform": f_text, "action": f_text,
       "priority": f_text, "severity": f_text, "message": f_text, "rationale": f_text,
       "category": f_text, "alert_type": f_text}


# ---------------------------------------------------------------- slides

def _cover(prs, period, kpis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    panel = slide.shapes.add_shape(1, 0, 0, W, H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = BRAND
    panel.line.fill.background()

    if LOGO_WHITE.exists():
        height = Inches(1.05)
        slide.shapes.add_picture(str(LOGO_WHITE), Inches(1.0), Inches(1.9),
                                 int(height * LOGO_RATIO), height)

    box = slide.shapes.add_textbox(Inches(1.0), Inches(3.2), W - Inches(2.0), Inches(2.6))
    frame = box.text_frame
    frame.word_wrap = True
    title = frame.paragraphs[0]
    title.text = "" if LOGO_WHITE.exists() else "Levista Foods"
    if title.runs:
        title.runs[0].font.size = Pt(52)
        title.runs[0].font.bold = True
        title.runs[0].font.color.rgb = WHITE
    for text, size in (("Marketing Performance Report", 30),
                       ("Amazon · Flipkart · Instamart · Zepto · BigBasket · Blinkit", 16)):
        para = frame.add_paragraph()
        para.text = text
        para.runs[0].font.size = Pt(size)
        para.runs[0].font.color.rgb = RGBColor(0xCF, 0xDA, 0xE2)
    start, end = period
    para = frame.add_paragraph()
    para.text = (f"{start:%d %B %Y} – {end:%d %B %Y}" if start and end
                 else "Reporting period not stated in source files")
    para.runs[0].font.size = Pt(18)
    para.runs[0].font.color.rgb = WHITE

    if kpis:
        strip = slide.shapes.add_textbox(Inches(1.0), Inches(5.6), W - Inches(2.0), Inches(0.8))
        _text(strip.text_frame,
              [f"{ai.inr(kpis['revenue'])} revenue   ·   {ai.inr(kpis['spend'])} spend   ·   "
               f"₹{kpis['roas']:.2f} returned per ₹1   ·   {int(kpis['orders']):,} orders"],
              size=16, color=WHITE)


def _executive(prs, engine, period, kpis):
    slide = _slide(prs, "Executive Summary",
                   "The whole picture in one view — every platform, one reporting period")
    # overall_kpis returns {} when a period has no primary rows — a real state on a
    # fresh install, and one that must not take the whole deck down with a KeyError.
    if not kpis:
        box = slide.shapes.add_textbox(MARGIN, Inches(2.0), W - 2 * MARGIN, Inches(1.2))
        _text(box.text_frame, [
            "No data loaded for this period.",
            "Upload the platform exports on the Data & Uploads page, then regenerate.",
        ], size=16, color=RED)
        return
    top = _kpi_cards(slide, [
        ("Revenue", ai.inr(kpis["revenue"]), GREEN),
        ("Ad Spend", ai.inr(kpis["spend"]), BRAND),
        ("ROAS", f"₹{kpis['roas']:.2f}", GREEN if kpis["roas"] >= config.HEALTHY_ROAS else AMBER),
        ("Orders", f"{int(kpis['orders']):,}", BRAND),
        ("Return on Investment", ai.pct(kpis["roi"]), GREEN if kpis["roi"] > 0 else RED),
    ])
    table = metrics.platform_comparison(engine)
    _table(slide, table,
           [("Platform", "platform"), ("Revenue", "revenue"), ("Spend", "spend"),
            ("Orders", "orders"), ("ROAS", "roas"), ("Share of revenue", "revenue_share")],
           MARGIN, top, Inches(7.2), Inches(3.4), FMT)

    box = slide.shapes.add_textbox(Inches(8.0), top, Inches(4.8), Inches(3.6))
    if table.empty:
        _text(box.text_frame, ["Headlines", "No platform data was loaded for this period."],
              size=14, bold_first=True, color=BRAND)
        return
    best = table.iloc[0]
    worst = table.sort_values("roas").iloc[0]
    # With a single platform loaded, best and worst are the same row — calling it both
    # the biggest earner and "the first place to act" reads as a contradiction.
    if len(table) == 1:
        _text(box.text_frame, [
            "Headlines",
            f"• {best['platform']} is the only platform loaded this period: "
            f"{ai.inr(best['revenue'])} of revenue on {ai.inr(best['spend'])} of spend.",
            f"• That is ₹{best['roas']:.2f} back for every ₹1 spent.",
            "• Load the other platforms' exports to compare them against each other.",
            "• ROAS means Revenue ÷ Ad Spend. Above ₹3 is healthy, below ₹1 loses money.",
        ], size=14, bold_first=True, color=BRAND)
        return
    _text(box.text_frame, [
        "Headlines",
        f"• {best['platform']} is the biggest revenue source at {ai.inr(best['revenue'])} "
        f"({ai.pct(best['revenue_share'])} of the total).",
        f"• {worst['platform']} returns only ₹{worst['roas']:.2f} per ₹1 spent — the weakest "
        "use of budget and the first place to act.",
        f"• Across all platforms every ₹1 of spend returned ₹{kpis['roas']:.2f}.",
        "• ROAS means Revenue ÷ Ad Spend. Above ₹3 is healthy, below ₹1 loses money.",
    ], size=14, bold_first=True, color=BRAND)


def _revenue_overview(prs, engine, period):
    slide = _slide(prs, "Revenue Overview", "Where the sales came from")
    table = metrics.platform_comparison(engine)
    _chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, table["platform"],
           {"Revenue": table["revenue"]}, MARGIN, Inches(1.5), Inches(7.0), Inches(4.4),
           title="Revenue by platform", legend=False)
    _chart(slide, XL_CHART_TYPE.DOUGHNUT, table["platform"],
           {"Share": table["revenue"]}, Inches(7.9), Inches(1.5), Inches(4.9), Inches(4.4),
           title="Share of total revenue")
    note = slide.shapes.add_textbox(MARGIN, Inches(6.1), W - 2 * MARGIN, Inches(0.9))
    # One platform is a normal state — a month can be loaded a file at a time, and
    # early on only one platform's export exists. Reaching for a second row here
    # crashed the whole rebuild ("single positional indexer is out-of-bounds"), which
    # surfaced to the user as uploads silently not working.
    top = table.head(2)
    if top.empty:
        message = "No platform revenue was loaded for this period."
    elif len(top) == 1:
        row = top.iloc[0]
        message = (f"{row['platform']} is the only platform loaded this period, so it "
                   f"accounts for all {ai.inr(row['revenue'])} of advertising revenue.")
    else:
        message = (f"{top.iloc[0]['platform']} and {top.iloc[1]['platform']} together account "
                   f"for {ai.pct(top['revenue_share'].sum())} of all advertising revenue "
                   "this period.")
    _text(note.text_frame, [message], size=14, color=BRAND)


def _spend_overview(prs, engine, period):
    slide = _slide(prs, "Spend Overview", "Where the budget went, and what it returned")
    table = metrics.platform_comparison(engine)
    _chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, table["platform"],
           {"Spend": table["spend"], "Revenue": table["revenue"]},
           MARGIN, Inches(1.5), Inches(7.6), Inches(4.4),
           title="Spend against revenue, by platform")
    _chart(slide, XL_CHART_TYPE.BAR_CLUSTERED, table["platform"], {"ROAS": table["roas"]},
           Inches(8.4), Inches(1.5), Inches(4.4), Inches(4.4),
           title="Return per ₹1 spent", legend=False)
    losing = table[table["roas"] < config.BREAKEVEN_ROAS]
    note = slide.shapes.add_textbox(MARGIN, Inches(6.1), W - 2 * MARGIN, Inches(0.9))
    message = ("Every platform returned more than it cost this period."
               if losing.empty else
               f"{', '.join(losing['platform'])} spent more than it earned back — "
               f"{ai.inr(losing['spend'].sum())} of budget returned "
               f"{ai.inr(losing['revenue'].sum())}.")
    _text(note.text_frame, [message], size=14, color=RED if not losing.empty else GREEN)


def _platform_comparison(prs, engine, period):
    slide = _slide(prs, "Platform Comparison", "Ranked by revenue, with efficiency alongside")
    table = metrics.platform_comparison(engine)
    _table(slide, table,
           [("Platform", "platform"), ("Revenue", "revenue"), ("Spend", "spend"),
            ("Orders", "orders"), ("ROAS", "roas"), ("CTR", "ctr"),
            ("Conversion", "conv_rate"), ("Revenue share", "revenue_share")],
           MARGIN, Inches(1.5), W - 2 * MARGIN, Inches(3.0), FMT)
    box = slide.shapes.add_textbox(MARGIN, Inches(5.0), W - 2 * MARGIN, Inches(2.0))
    _text(box.text_frame, [
        "How to read this",
        "• ROAS is the headline efficiency number: revenue divided by spend.",
        "• A blank CTR or conversion figure means that platform's export does not report "
        "clicks for most of its spend, so the number would not be comparable.",
        "• Revenue share shows how dependent the business is on each platform.",
    ], size=13, bold_first=True, color=BRAND)


def _platform_slide(prs, engine, platform, period):
    slide = _slide(prs, f"{platform} Analysis", f"Performance, causes and next steps for {platform}")
    table = metrics.platform_comparison(engine)
    row = table[table["platform"] == platform]
    if row.empty:
        box = slide.shapes.add_textbox(MARGIN, Inches(2.0), W - 2 * MARGIN, Inches(1.0))
        _text(box.text_frame,
              [f"No {platform} data was found in this month's input files."], size=16, color=RED)
        return
    row = row.iloc[0]
    top = _kpi_cards(slide, [
        ("Revenue", ai.inr(row["revenue"]), GREEN),
        ("Spend", ai.inr(row["spend"]), BRAND),
        ("ROAS", f"₹{row['roas']:.2f}",
         GREEN if row["roas"] >= config.HEALTHY_ROAS
         else AMBER if row["roas"] >= config.BREAKEVEN_ROAS else RED),
        ("Orders", f"{int(row['orders'] or 0):,}", BRAND),
    ])

    campaigns = metrics.collapse(metrics.campaigns(engine, platform=platform), "campaign")
    if not campaigns.empty:
        _table(slide, campaigns.head(6),
               [("Top campaigns", "campaign_name"), ("Spend", "spend"),
                ("Revenue", "revenue"), ("ROAS", "roas")],
               MARGIN, top, Inches(7.0), Inches(2.6), FMT)

    narrative = ai.load(engine, insights_table)
    story = narrative[narrative["platform"] == platform]
    box = slide.shapes.add_textbox(Inches(7.7), top, Inches(5.1), Inches(4.2))
    if story.empty:
        _text(box.text_frame, ["No narrative available."], size=13)
    else:
        item = story.iloc[0]
        _text(box.text_frame, [
            "What happened", item["what_happened"],
            "Why", item["why_it_happened"],
            "What to do next", item["what_to_do_next"],
        ], size=12, color=GREY)
        for i, para in enumerate(box.text_frame.paragraphs):
            if i % 2 == 0:
                for run in para.runs:
                    run.font.bold = True
                    run.font.color.rgb = BRAND


def _entity_slide(prs, engine, title, subtitle, df, columns, chart_label, chart_value="revenue"):
    slide = _slide(prs, title, subtitle)
    if df.empty:
        box = slide.shapes.add_textbox(MARGIN, Inches(2.0), W - 2 * MARGIN, Inches(1.0))
        _text(box.text_frame, ["No source data for this view in this month's files."],
              size=16, color=RED)
        return
    top = df.head(10)
    _table(slide, top, columns, MARGIN, Inches(1.5), Inches(7.4), Inches(4.6), FMT)
    labels = [str(v)[:22] for v in top[chart_label]]
    _chart(slide, XL_CHART_TYPE.BAR_CLUSTERED, labels[::-1],
           {chart_value.title(): list(top[chart_value])[::-1]},
           Inches(8.1), Inches(1.5), Inches(4.7), Inches(4.6),
           title=f"{chart_value.title()} — top 10", legend=False)


def _campaign_extremes(prs, engine, best: bool):
    title = "Top Performing Campaigns" if best else "Underperforming Campaigns"
    subtitle = ("Highest return on every rupee — protect and scale these"
                if best else "Spending more than they return — act on these first")
    slide = _slide(prs, title, subtitle)
    df = metrics.collapse(metrics.campaigns(engine), "campaign", ["platform"])
    df = df[df["spend"].fillna(0) >= config.MIN_SPEND_FOR_ALERT]
    if df.empty:
        box = slide.shapes.add_textbox(MARGIN, Inches(2.0), W - 2 * MARGIN, Inches(1.0))
        _text(box.text_frame, ["No campaign data available."], size=16, color=RED)
        return
    df = df.sort_values("roas", ascending=not best, na_position="last").head(10)
    _table(slide, df,
           [("Platform", "platform"), ("Campaign", "campaign_name"), ("Spend", "spend"),
            ("Revenue", "revenue"), ("Orders", "orders"), ("ROAS", "roas")],
           MARGIN, Inches(1.5), W - 2 * MARGIN, Inches(4.4), FMT)
    box = slide.shapes.add_textbox(MARGIN, Inches(6.2), W - 2 * MARGIN, Inches(0.9))
    if best:
        _text(box.text_frame, [
            f"These ten campaigns turned {ai.inr(df['spend'].sum())} into "
            f"{ai.inr(df['revenue'].sum())}. Increasing budget here is the lowest-risk growth "
            "available."], size=14, color=GREEN)
    else:
        loss = df["spend"].sum() - df["revenue"].sum()
        _text(box.text_frame, [
            f"These ten campaigns spent {ai.inr(df['spend'].sum())} and returned "
            f"{ai.inr(df['revenue'].sum())} — a shortfall of {ai.inr(loss)}. "
            "Pause or re-bid them before adding budget anywhere else."], size=14, color=RED)


def _budget_wastage(prs, engine):
    slide = _slide(prs, "Budget Wastage Analysis",
                   "Money spent that produced no sales at all")
    rows = []
    for entity, table in (("Campaigns", metrics.campaigns(engine)),
                          ("Products", metrics.products(engine)),
                          ("Keywords", metrics.keywords(engine)),
                          ("Cities", metrics.cities(engine))):
        leak = metrics.wasted_spend(table, config.MIN_SPEND_FOR_ALERT)
        rows.append({"area": entity, "count": len(leak),
                     "spend": leak["spend"].sum() if not leak.empty else 0})
    summary = pd.DataFrame(rows)
    total = summary["spend"].sum()
    _kpi_cards(slide, [("Wasted spend identified", ai.inr(total), RED),
                       ("Items with spend but no sales",
                        f"{int(summary['count'].sum()):,}", RED)])
    _chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, summary["area"],
           {"Wasted spend": summary["spend"]},
           MARGIN, Inches(3.0), Inches(6.4), Inches(3.4),
           title="Where the waste sits", legend=False)
    worst = metrics.wasted_spend(metrics.collapse(metrics.keywords(engine), "keyword",
                                                  ["platform"]),
                                 config.MIN_SPEND_FOR_ALERT).head(8)
    if not worst.empty:
        _table(slide, worst,
               [("Platform", "platform"), ("Keyword", "keyword"), ("Spend", "spend"),
                ("Impressions", "impressions")],
               Inches(7.2), Inches(3.0), Inches(5.6), Inches(3.4), FMT)


def _opportunities(prs, engine):
    slide = _slide(prs, "Opportunities", "Where more budget would earn the most")
    recs = ai.load(engine, recommendations)
    grow = recs[recs["action"] == "Increase budget"].sort_values(
        "impact_value", ascending=False).head(10)
    if grow.empty:
        box = slide.shapes.add_textbox(MARGIN, Inches(2.0), W - 2 * MARGIN, Inches(1.0))
        _text(box.text_frame, ["No scaling opportunities were identified this period."],
              size=16, color=GREY)
        return
    _table(slide, grow,
           [("Platform", "platform"), ("Type", "entity_type"), ("Name", "entity_name"),
            ("Revenue today", "impact_value"), ("Why", "rationale")],
           MARGIN, Inches(1.5), W - 2 * MARGIN, Inches(4.6), FMT)
    box = slide.shapes.add_textbox(MARGIN, Inches(6.3), W - 2 * MARGIN, Inches(0.8))
    _text(box.text_frame, [
        "Each of these already earns well above break-even on a small budget. Raising their "
        "budgets is the cheapest growth available this month."], size=14, color=GREEN)


def _recommendations(prs, engine):
    slide = _slide(prs, "Recommendations", "Ranked by priority and money at stake")
    recs = ai.load(engine, recommendations)
    if recs.empty:
        box = slide.shapes.add_textbox(MARGIN, Inches(2.0), W - 2 * MARGIN, Inches(1.0))
        _text(box.text_frame, ["No recommendations were generated."], size=16, color=GREY)
        return
    recs["_order"] = recs["priority"].map(ai.PRIORITY_ORDER).fillna(9)
    top = recs.sort_values(["_order", "impact_value"], ascending=[True, False]).head(10)
    _table(slide, top,
           [("Priority", "priority"), ("Platform", "platform"), ("Action", "action"),
            ("Name", "entity_name"), ("Money at stake", "impact_value")],
           MARGIN, Inches(1.5), Inches(8.0), Inches(4.6), FMT)
    counts = recs.groupby("action").size().sort_values(ascending=False)
    _chart(slide, XL_CHART_TYPE.BAR_CLUSTERED, list(counts.index)[::-1],
           {"Actions": list(counts.values)[::-1]},
           Inches(8.7), Inches(1.5), Inches(4.1), Inches(4.6),
           title="Recommendations by action", legend=False)


def _action_plan(prs, engine):
    slide = _slide(prs, "Action Plan", "Who does what, in what order")
    recs = ai.load(engine, recommendations)
    table = metrics.platform_comparison(engine)
    losing = table[table["roas"] < config.BREAKEVEN_ROAS]["platform"].tolist()
    pause_spend = recs[recs["action"] == "Pause"]["impact_value"].sum()
    cut_spend = recs[recs["action"] == "Reduce spend"]["impact_value"].sum()
    grow_count = int((recs["action"] == "Increase budget").sum())

    steps = [
        ("This week — stop the losses",
         f"Pause the {int((recs['action'] == 'Pause').sum())} items spending with no sales at "
         f"all ({ai.inr(pause_spend)} of budget freed immediately)."),
        ("This week — cut the leaks",
         f"Reduce bids on the {int((recs['action'] == 'Reduce spend').sum())} items returning "
         f"below break-even (about {ai.inr(cut_spend)} currently lost)."),
        ("Next two weeks — move the money",
         f"Redirect the freed budget into the {grow_count} proven winners listed on the "
         "Opportunities slide."),
        ("Next two weeks — fix the weak platform",
         (f"{', '.join(losing)} is below break-even. Review pricing, product-page content and "
          "keyword match types before restoring budget.") if losing
         else "All platforms are above break-even — hold budgets and keep optimising."),
        ("This month — close the reporting gaps",
         "Pull the missing exports listed on the Alerts sheet so next month's report is complete."),
    ]
    top = Inches(1.5)
    for i, (heading, detail) in enumerate(steps):
        box = slide.shapes.add_shape(1, MARGIN, top + Inches(1.05) * i,
                                     W - 2 * MARGIN, Inches(0.92))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDC)
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.2)
        head = frame.paragraphs[0]
        head.text = heading
        head.runs[0].font.size = Pt(13)
        head.runs[0].font.bold = True
        head.runs[0].font.color.rgb = BRAND
        body = frame.add_paragraph()
        body.text = detail
        body.runs[0].font.size = Pt(12)
        body.runs[0].font.color.rgb = GREY


def _appendix(prs, engine, counts):
    slide = _slide(prs, "Appendix", "Data sources, definitions and known gaps")
    alert_df = ai.load(engine, alerts)
    gaps = alert_df[alert_df["alert_type"] == "Missing report"]
    lines = [
        "Where the numbers come from",
        f"• {counts['files']} export files were processed; {counts['duplicates']} were exact "
        "duplicates of another file and were counted once.",
        "• Platform totals use one report per platform so no rupee is counted twice.",
        "• CTR, CPC, CPM, ROAS and conversion rate are recalculated from impressions, clicks, "
        "spend, orders and revenue rather than copied, so platforms are comparable.",
        "",
        "Definitions",
        "• ROAS = Revenue ÷ Ad Spend.   • ROI = (Revenue − Spend) ÷ Spend.",
        "• CTR = Clicks ÷ Impressions.  • Conversion rate = Orders ÷ Clicks.",
        "",
        "Known gaps this period",
    ]
    lines += ([f"• {row['message']}" for _, row in gaps.iterrows()]
              or ["• None — every expected report was supplied."])
    box = slide.shapes.add_textbox(MARGIN, Inches(1.4), W - 2 * MARGIN, Inches(5.6))
    _text(box.text_frame, lines, size=12, color=GREY)
    for para in box.text_frame.paragraphs:
        if para.text and not para.text.startswith("•"):
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = BRAND


# ---------------------------------------------------------------- entry point

def build(engine=None, path=None, period_label=None):
    engine = engine or get_engine()
    period_label = period_label or metrics.latest_period(engine)
    period = metrics.reporting_period(engine, period_label)
    kpis = metrics.overall_kpis(engine, period_label)

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    _cover(prs, period, kpis)                                            # 1
    _executive(prs, engine, period, kpis)                                # 2
    _revenue_overview(prs, engine, period)                               # 3
    _spend_overview(prs, engine, period)                                 # 4
    _platform_comparison(prs, engine, period)                            # 5
    for platform in ["Amazon", "Flipkart", "Instamart", "BigBasket", "Zepto", "Blinkit"]:
        _platform_slide(prs, engine, platform, period)                   # 6-11

    _entity_slide(prs, engine, "Product Performance",                    # 12
                  "Best selling advertised products across all platforms",
                  metrics.collapse(metrics.products(engine), "product", ["platform"]),
                  [("Platform", "platform"), ("Product", "product_name"),
                   ("Spend", "spend"), ("Revenue", "revenue"), ("ROAS", "roas")],
                  "product_name")
    _entity_slide(prs, engine, "Keyword Performance",                    # 13
                  "The search terms driving the most revenue",
                  metrics.collapse(metrics.keywords(engine), "keyword", ["platform"]),
                  [("Platform", "platform"), ("Keyword", "keyword"), ("Spend", "spend"),
                   ("Revenue", "revenue"), ("ROAS", "roas")], "keyword")
    _entity_slide(prs, engine, "City Performance",                       # 14
                  "Where in India the sales are coming from",
                  metrics.collapse(metrics.cities(engine), "city", ["platform"]),
                  [("Platform", "platform"), ("City", "city"), ("Spend", "spend"),
                   ("Revenue", "revenue"), ("ROAS", "roas")], "city")

    _campaign_extremes(prs, engine, best=True)                           # 15
    _campaign_extremes(prs, engine, best=False)                          # 16
    _budget_wastage(prs, engine)                                         # 17
    _opportunities(prs, engine)                                          # 18
    _recommendations(prs, engine)                                        # 19
    _action_plan(prs, engine)                                            # 20

    with engine.connect() as conn:
        from db.models import uploaded_files
        files = pd.read_sql(uploaded_files.select(), conn)
    _appendix(prs, engine, {                                             # 21
        "files": len(files),
        "duplicates": int((files["processing_status"] == "duplicate").sum()),
    })

    path = path or (config.OUTPUT_DIR / "Levista_Performance_Report.pptx")
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    out, slides = build()
    print(f"Wrote {out} ({slides} slides)")
